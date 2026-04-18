import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── カラーパレット ──
C_BG     = RGBColor(0x1A, 0x25, 0x2F)
C_ACCENT = RGBColor(0xF0, 0xC0, 0x27)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xCC, 0xDD, 0xEE)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_BLUE   = RGBColor(0x29, 0x80, 0xB9)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_GRAY   = RGBColor(0xAA, 0xBB, 0xCC)
C_DARK2  = RGBColor(0x22, 0x33, 0x44)
C_DARK3  = RGBColor(0x20, 0x30, 0x40)

BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG
    return slide


def txb(slide, text, left, top, width, height,
        size=24, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tf = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def section_header(slide, title):
    txb(slide, title, 0.5, 0.3, 12, 0.6, size=14, color=C_ACCENT, bold=True)
    rect(slide, 0.5, 0.85, 12.3, 0.04, C_ACCENT)


# ══════════════════════════════════════
# スライド1: タイトル
# ══════════════════════════════════════
s = add_slide()
rect(s, 0, 3.05, 13.33, 0.07, C_ACCENT)
txb(s, '競馬AIを自作してみた', 1, 1.1, 11.3, 1.7,
    size=54, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txb(s, '〜 Python × 機械学習 × 実データ 〜', 1, 3.3, 11.3, 0.7,
    size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
txb(s, 'JRA-VANデータ  ×  LightGBM  ×  LambdaMART', 1, 4.15, 11.3, 0.6,
    size=18, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
txb(s, '2026.3', 10.5, 6.8, 2.5, 0.5, size=14, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════
# スライド2: 一言で言うと
# ══════════════════════════════════════
s = add_slide()
section_header(s, '一言で言うと')
txb(s, '出馬表を渡すだけで', 1, 1.4, 11.3, 1.0,
    size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txb(s, '自動で予想PDFが出てくる仕組みを作った', 1, 2.3, 11.3, 1.0,
    size=34, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

boxes = [
    ('📄 出馬表CSV', C_BLUE,   0.5),
    ('🤖 AI予測',   C_GREEN,  4.6),
    ('📑 PDF出力',  C_ORANGE, 8.7),
]
for txt, col, x in boxes:
    rect(s, x, 4.0, 3.7, 1.4, C_DARK2)
    txb(s, txt, x, 4.1, 3.7, 1.2, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

txb(s, '→', 4.3, 4.35, 0.5, 0.7, size=28, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
txb(s, '→', 8.4, 4.35, 0.5, 0.7, size=28, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
txb(s, '毎週全自動で実行', 4.0, 5.65, 5.3, 0.6,
    size=16, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# スライド3: データ規模
# ══════════════════════════════════════
s = add_slide()
section_header(s, '使ったデータ')

stats = [
    ('約20年分',  'JRA-VANレースデータ（2005〜2026）', C_BLUE),
    ('約66万頭',  '学習データ（2013〜2025年）',        C_GREEN),
    ('272列',    '1頭あたりの特徴量数',                 C_ORANGE),
    ('約1.3GB',  '特徴量CSV（parquet圧縮で219MB）',    C_ACCENT),
]
for i, (num, desc, col) in enumerate(stats):
    x = 0.5 + (i % 2) * 6.4
    y = 1.2 + (i // 2) * 2.7
    rect(s, x, y, 5.9, 2.3, C_DARK2)
    rect(s, x, y, 5.9, 0.08, col)
    txb(s, num, x + 0.2, y + 0.2, 5.5, 1.0, size=32, bold=True, color=col)
    txb(s, desc, x + 0.2, y + 1.2, 5.5, 0.8, size=15, color=C_LIGHT)

# ══════════════════════════════════════
# スライド4: 技術スタック
# ══════════════════════════════════════
s = add_slide()
section_header(s, '使った技術')

techs = [
    ('Python',            '全体の基盤・データ処理・自動化',              C_BLUE),
    ('pandas / numpy',    'データ処理・特徴量エンジニアリング',           C_GREEN),
    ('LightGBM',          '勝率予測（勾配ブースティング・分類器）',        C_ORANGE),
    ('LambdaMART',        'レース内順位予測（学習ランキング）',            C_RED),
    ('python-pptx',       'このスライド自体をPythonで自動生成',           C_ACCENT),
    ('Chrome Headless',   'HTML → PDF自動変換（毎週の予想出力）',        C_GRAY),
]
for i, (name, desc, col) in enumerate(techs):
    y = 1.1 + i * 0.95
    rect(s, 0.5, y, 12.3, 0.82, C_DARK2)
    rect(s, 0.5, y, 0.12, 0.82, col)
    txb(s, name, 0.75, y + 0.1, 3.8, 0.6, size=17, bold=True, color=col)
    txb(s, desc, 4.8, y + 0.12, 7.8, 0.6, size=15, color=C_LIGHT)

# ══════════════════════════════════════
# スライド5: AIモデルの仕組み
# ══════════════════════════════════════
s = add_slide()
section_header(s, 'AIモデルの仕組み')

rect(s, 0.5, 1.1, 5.8, 2.6, RGBColor(0x1A, 0x3A, 0x2A))
rect(s, 0.5, 1.1, 5.8, 0.5, C_GREEN)
txb(s, 'Step 1：LightGBM（分類）', 0.6, 1.13, 5.6, 0.45,
    size=15, bold=True, color=RGBColor(0x1A, 0x25, 0x2F))
txb(s, '各馬の「単勝確率」を計算\n\n騎手勝率・調教師勝率・距離適性・\n近走成績など272特徴量を入力', 0.7, 1.7, 5.5, 1.8, size=14, color=C_WHITE)

txb(s, '→  両方の結果を統合', 5.8, 2.2, 1.7, 0.65, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

rect(s, 7.2, 1.1, 5.6, 2.6, RGBColor(0x1A, 0x28, 0x3A))
rect(s, 7.2, 1.1, 5.6, 0.5, C_BLUE)
txb(s, 'Step 2：LambdaMART（ランキング）', 7.3, 1.13, 5.4, 0.45,
    size=15, bold=True, color=RGBColor(0x1A, 0x25, 0x2F))
txb(s, 'レース内の「着順」を最適化\n\n同じレースの馬同士を正しく\n順位付けするように学習', 7.3, 1.7, 5.3, 1.8, size=14, color=C_WHITE)

txb(s, '2種類のモデルを並走させる', 0.5, 3.9, 12.3, 0.6,
    size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

rect(s, 0.5, 4.6, 5.8, 2.5, RGBColor(0x1A, 0x3A, 0x2A))
txb(s, '🟢 距離モデル', 0.7, 4.7, 5.4, 0.6, size=17, bold=True, color=C_GREEN)
txb(s, '会場 × 距離 のキーで学習\n例: 中山ダート1200m専用モデル\n（コース適性を重視）', 0.7, 5.25, 5.4, 1.6, size=14, color=C_WHITE)

rect(s, 7.2, 4.6, 5.6, 2.5, RGBColor(0x1A, 0x28, 0x3A))
txb(s, '🔵 クラスモデル', 7.4, 4.7, 5.2, 0.6, size=17, bold=True, color=C_BLUE)
txb(s, '芝ダ × 距離帯 × クラス のキーで学習\n例: ダート短距離・未勝利専用モデル\n（クラス適性を重視）', 7.4, 5.25, 5.2, 1.6, size=14, color=C_WHITE)

txb(s, '両モデルでランカー1位 → 本命候補', 0.5, 7.1, 12.3, 0.4,
    size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# スライド6: 印システム
# ══════════════════════════════════════
s = add_slide()
section_header(s, '印（マーク）システム')

marks = [
    ('激熱', '両モデルでランカー1位\ncombo_gap≥15\n偏差値差≥10\nオッズ3倍以上', RGBColor(0xB7, 0x41, 0x0E), RGBColor(0x2A, 0x18, 0x10)),
    ('◎',   '両モデルでランカー1位\ncombo_gap≥10\nオッズ3倍以上',              RGBColor(0x1A, 0x6E, 0x3C), RGBColor(0x18, 0x28, 0x1C)),
    ('〇',   '両モデルでランカー1位\ncombo_gap＜10\nオッズ3倍以上',             RGBColor(0x27, 0xAE, 0x60), RGBColor(0x18, 0x28, 0x1C)),
    ('☆',   '片方が2〜3位・もう片方も\n3位以内に入っている\nオッズ5倍以上',      RGBColor(0x1F, 0x61, 0x8D), RGBColor(0x10, 0x1C, 0x2A)),
]
for i, (mark, cond, col, bg) in enumerate(marks):
    x = 0.5 + i * 3.2
    rect(s, x, 1.1, 3.0, 5.9, bg)
    rect(s, x, 1.1, 3.0, 1.1, col)
    txb(s, mark, x, 1.1, 3.0, 1.1, size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txb(s, cond, x + 0.15, 2.35, 2.7, 3.5, size=14, color=C_LIGHT)

txb(s, '両モデル一致が鍵。片方だけでは買わない。', 0.5, 7.1, 12.3, 0.4,
    size=15, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# スライド7: ROI検証結果
# ══════════════════════════════════════
s = add_slide()
section_header(s, '検証結果（2026年 実データ out-of-sample）')
txb(s, '学習データと完全に切り分けた2026年のリアルデータで検証',
    0.5, 1.0, 12.3, 0.5, size=14, color=C_GRAY, italic=True)

rect(s, 0.5, 1.55, 12.3, 0.5, RGBColor(0x2C, 0x3E, 0x50))
headers = [('戦略条件', 6.0, 0.5), ('頭数', 6.5, 1.5), ('単勝ROI', 8.0, 2.0)]
for hdr, cx, cw in headers:
    txb(s, hdr, cx + 0.1, 1.6, cw, 0.4, size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

rows = [
    ('両Rnk=1（オッズ制限なし）',                           '339頭', '+21.4%',  C_LIGHT, False),
    ('両Rnk=1  &  オッズ3倍以上',                          '169頭', '+52.7%',  C_LIGHT, False),
    ('両Rnk=1  &  odds≥3  &  偏差値差≥10',                ' 62頭', '+127.3%', C_GREEN, True),
    ('穴馬: combo_rank=2  &  両Rnk≤3  &  odds≥5  &  偏差値差≥10', '29頭', '+183.4%', C_ACCENT, True),
    ('旧ルール（odds≥5  &  偏差値差≥20）',                   '22頭', '-71.4%',  C_RED,   False),
]
for i, (strat, n, roi, col, hl) in enumerate(rows):
    y = 2.15 + i * 0.9
    bg = RGBColor(0x25, 0x38, 0x25) if hl else C_DARK2
    rect(s, 0.5, y, 12.3, 0.78, bg)
    if hl:
        rect(s, 0.5, y, 0.15, 0.78, col)
    txb(s, strat, 0.75, y + 0.1, 5.6, 0.58, size=14, bold=hl, color=col)
    txb(s, n,    6.6, y + 0.1, 1.3, 0.58, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txb(s, roi,  8.1, y + 0.08, 1.9, 0.62, size=16, bold=hl, color=col, align=PP_ALIGN.CENTER)

txb(s, '※ JRAの控除率は約25%。ROI+21%でも競馬場平均の期待値の約2倍',
    0.5, 7.05, 12.3, 0.42, size=12, color=C_GRAY, italic=True)

# ══════════════════════════════════════
# スライド8: こだわりポイント
# ══════════════════════════════════════
s = add_slide()
section_header(s, 'こだわりポイント')

points = [
    ('時系列リーク対策',
     '学習データ（〜2025年末）と検証データ（2026年〜）を完全分離。\n未来のデータが学習に混入しないよう厳格に管理。',
     C_RED),
    ('全自動パイプライン',
     '出馬表CSVを渡すだけで特徴量生成 → モデル予測 → HTML/PDF出力まで全自動。\n実運用で毎週使えるシステム。',
     C_GREEN),
    ('データ品質アラート',
     '地方競馬前走・海外帰り馬・特徴量欠損を自動検出してPDFに警告表示。\n最終判断は人間が行うハイブリッド運用。',
     C_ORANGE),
    ('2モデル並走 & 多角検証',
     '距離特化モデル・クラス特化モデルの2つを常に並走。\n両方が一致した場合のみ推奨とすることで精度向上。',
     C_BLUE),
]
for i, (title, body, col) in enumerate(points):
    x = 0.5 + (i % 2) * 6.4
    y = 1.1 + (i // 2) * 2.95
    rect(s, x, y, 5.9, 2.65, C_DARK3)
    rect(s, x, y, 0.13, 2.65, col)
    txb(s, title, x + 0.3, y + 0.15, 5.4, 0.65, size=16, bold=True, color=col)
    txb(s, body,  x + 0.3, y + 0.8,  5.4, 1.7,  size=13, color=C_LIGHT)

# ══════════════════════════════════════
# スライド9: 開発規模・開発者について
# ══════════════════════════════════════
s = add_slide()
section_header(s, '開発規模')

nums = [
    ('20本+',  'Pythonスクリプト',        C_GREEN),
    ('272列',  '1頭あたり特徴量数',        C_BLUE),
    ('約66万頭', '学習データ規模',          C_ORANGE),
    ('100本+', '学習済みモデル数\n（コース×クラス別）', C_ACCENT),
]
for i, (num, label, col) in enumerate(nums):
    x = 0.5 + i * 3.2
    rect(s, x, 1.1, 3.0, 2.3, C_DARK2)
    rect(s, x, 1.1, 3.0, 0.08, col)
    txb(s, num,   x, 1.25, 3.0, 1.1, size=30, bold=True, color=col,   align=PP_ALIGN.CENTER)
    txb(s, label, x, 2.3,  3.0, 0.9, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

txb(s, '開発者について', 0.5, 3.65, 12.3, 0.55, size=18, bold=True, color=C_ACCENT)
rect(s, 0.5, 4.15, 12.3, 0.04, RGBColor(0x44, 0x55, 0x66))

dev_points = [
    'プログラミング歴は長くない（開発初心者〜中級者）',
    '機械学習の専門家ではない',
    'それでも実用レベルのAIシステムを構築・毎週運用できている',
    'Claude AIとの対話を活用しながら設計・デバッグ・改善を継続中',
]
for i, pt in enumerate(dev_points):
    txb(s, '✅  ' + pt, 0.8, 4.35 + i * 0.73, 12.0, 0.65, size=16, color=C_WHITE)

# ══════════════════════════════════════
# スライド10: まとめ
# ══════════════════════════════════════
s = add_slide()
rect(s, 0, 2.85, 13.33, 0.07, C_ACCENT)
txb(s, 'まとめ', 0.5, 0.3, 12, 0.6, size=14, color=C_ACCENT, bold=True)
rect(s, 0.5, 0.85, 12.3, 0.04, C_ACCENT)

sums = [
    ('📄→📑', 'CSVを渡すだけで予想PDFが出る実用システムを構築',           C_WHITE),
    ('🧠',    'LightGBM + LambdaMARTの2段階AIを自作',               C_GREEN),
    ('📈',    '2026年リアルデータで単勝ROI +127% を記録',              C_ACCENT),
    ('🛡️',   'データ品質チェック・時系列リーク対策まで実装',               C_BLUE),
    ('🔧',    '初心者でも「動き続ける実用システム」を作れる時代になった',    C_ORANGE),
]
for i, (icon, text, col) in enumerate(sums):
    txb(s, icon + '  ' + text, 1.0, 1.1 + i * 1.05, 11.3, 0.9,
        size=19, bold=(col != C_WHITE), color=col)

txb(s, '「作ること」より「動き続けること」にこだわった',
    1.0, 6.6, 11.3, 0.7, size=16, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)


out = r'G:\マイドライブ\horse_racing_ai\競馬AI_プレゼン.pptx'
prs.save(out)
print('保存完了:', out)
